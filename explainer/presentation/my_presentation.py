from dataclasses import dataclass, field
from collections import namedtuple
from typing import Dict
from pptx import Presentation

from .my_slide import MySlide


@dataclass
class MyPresentation:
    """
    Represents a PowerPoint presentation.

    Attributes:
        filepath (str): The file path of the PowerPoint presentation.
        explanations (namedtuple): Named tuple to store slide explanations. Defaults to None.
        slides (Dict[int, MySlide]): Dictionary to store slide objects. Defaults to None.
    """

    filepath: str
    explanations: namedtuple = None
    slides: Dict[int, MySlide] = field(default_factory=dict)

    def parse(self):
        """
        Parses the PowerPoint presentation to extract slide data.

        Creates instances of the MySlide class for each slide and stores them in the self.slides dictionary.
        """
        try:
            prs = Presentation(f"{self.filepath}")
        except Exception as e:
            print(f"Error: {e}")
            return

        for slide_number, slide in enumerate(prs.slides, start=1):
            slide_obj = prs.slides[slide_number - 1]  # Access slide object using slide number - 1

            # Extract text from text boxes on the slide
            slide_text_boxes = [
                run.text.strip()
                for shape in slide_obj.shapes
                if shape.has_text_frame
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
                if run.text.strip()
            ]

            # Ignore the slide if there is no text
            if not slide_text_boxes:
                continue

            # Create a MySlide object and store it in the slides dictionary
            self.slides[slide_number] = MySlide(slide_number, slide_text_boxes)
